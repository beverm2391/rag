from dotenv import load_dotenv
import os
from pydantic import BaseModel, ValidationError
from pypdf import PdfReader
import re
from pptx import Presentation
import tiktoken
from functools import wraps
import inspect
import cohere
import anthropic
import openai
import hashlib
import uuid
import json
import requests
from typing import Type

from lib.model_config import MODELS
from lib.config import logger, load_env

# ! Text Processing ========================
def count_tokens(text: str, model: str = "gpt-4") -> int:
    try:
        encoding = tiktoken.encoding_for_model(model)
    except Exception as e:
        logger.debug(f"Model {model} not found in tiktoken. Defaulting to base model. Error: {e}")
        encoding = tiktoken.get_encoding("cl100k_base")
    return len(encoding.encode(text)) if text else 0

def pdf_to_text(fpath: str): return''.join([page.extract_text() for page in PdfReader(fpath).pages]).replace('\n', ' ').strip()

def clean_text(text):
    text = re.sub(r'\n\s*\n', '\n', text)
    return re.sub(r' {2,}', ' ', text).replace("\n", " ").strip()

def PPTX_to_text(path: str):
    """Extract the text from a PPTX file."""
    if not os.path.exists(path): raise FileNotFoundError(f"File {path} not found")
    
    prs = Presentation(path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
    return "\n".join(text_runs)

# ! Asyncio Nesting ========================

def auto_nest_asyncio(func):
    import nest_asyncio

    @wraps(func)
    def wrapper(*args, **kwargs):
        try:
            return func(*args, **kwargs)
        except RuntimeError as e:
            print(f"Caught {e}", "\nThis is usually becasue a jupyter notebook has a running event loop. No problem - automatically nesting event loops...")
            if str(e) == 'asyncio.run() cannot be called from a running event loop':
                nest_asyncio.apply()
                return func(*args, **kwargs)
            else:
                raise e
    return wrapper


# ! Stream Validation ========================

def validate_stream(res):
    """Helper function to validate a generator."""
    assert res is not None, "Chat response is None"
    assert inspect.isgenerator(res), "Chat response is not a generator"

    while True:
        try:
            next(res)
        except StopIteration:
            break
        except Exception as e:
            print(e)
            assert False, "Error in stream"
        finally:
            pass


# ! Server API Key Generation ========================

def gen_api_key():
    data = str(uuid.uuid4()).encode('utf-8')
    return hashlib.sha256(data).hexdigest()


# ! HTTP Requests ========================

def post(URL: str, headers: dict = {}, data: dict = {}): return requests.post(URL, headers=headers, json=data)
def get(URL: str, headers: dict = {}, params: dict = {}): return requests.get(URL, headers=headers, params=params)

# ! List Available Models ========================
class Models:
    _models = MODELS
    _env = load_env()

    @staticmethod
    def cohere():
        """Return a list of all models from the Cohere API."""
        client = cohere.Client(api_key=Models._env['COHERE_API_KEY'])
        return [model.__dict__ for model in client.models.list().models]

    # TODO - add Anthropics models (couldn't find a method to list models in the API docs)
    # ? DOCS: https://docs.anthropic.com/claude/docs/models-overview
    @staticmethod
    def anthropic():
        pass

    @staticmethod
    def openai():
        client = openai.Client(api_key=Models._env['OPENAI_API_KEY'])
        return [model.__dict__ for model in client.models.list()]
    
# ! Stream Serialization ========================

def stream_dicts_as_json(dict_stream):
    """
    A generator that takes a stream of dictionaries and yields each
    as a serialized JSON string suitable for SSE (Server-Sent Events).

    Args:
    dict_stream (Iterator): An iterator that yields dictionaries.

    Yields:
    str: A serialized JSON string for each dictionary, formatted for SSE.
    """
    for dict_item in dict_stream:
        json_str = json.dumps(dict_item) # Convert the dictionary to a JSON string
        yield f"data: {json_str}\n\n" # Format for SSE protocol and yield


# ! Cost Calculation ========================

def calculate_cost(model: str, input_tokens: int, output_tokens: int):
    if not model in MODELS: raise ValueError(f"Model {model} not found in MODELS")
    price_obj = MODELS[model].get('price_per_million_tokens', None)
    if not price_obj: return None
    
    input_tokens_cost = input_tokens * price_obj['input'] / 1000000
    output_tokens_cost = output_tokens * price_obj['output'] / 1000000
    
    # ? I'm intentionally keeping precision (not rounding) so that it's precise for the frontend/user
    return {
        "input_tokens_cost": input_tokens_cost,
        "output_tokens_cost": output_tokens_cost,
        "total_cost": input_tokens_cost + output_tokens_cost,
    }

def verbose_validate(model: Type[BaseModel], data: dict) -> None:
    try:
        return model(**data)
    except ValidationError as e:
        # error_fields = {err['loc'][0] for err in e.errors()}
        model_fields = set(model.model_fields.keys())
        data_fields = set(data.keys())

        missing_in_data = model_fields - data_fields
        extra_in_data = data_fields - model_fields

        if missing_in_data:
            print(f"Fields missing in data: {missing_in_data}")
        if extra_in_data:
            print(f"Extra fields in data: {extra_in_data}")

        print("Validation errors:")
        for error in e.errors():
            loc = " -> ".join(map(str, error["loc"]))
            msg = error["msg"]
            print(f"{loc}: {msg}")

        raise e